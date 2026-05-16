# ============================================================
#  file_editor.py — Editor Universal de Arquivos para o R.O.N.Y
#
#  Suporte completo a:
#    • TXT / MD / CSV / JSON / código (qualquer texto)
#    • Word  (.docx)   — via python-docx
#    • Excel (.xlsx)   — via openpyxl
#    • PowerPoint (.pptx) — via python-pptx
#    • HTML dashboards — geração automática com Chart.js
#
#  Fluxo de uso:
#    1. Usuário faz pedido por voz / texto
#    2. editar_arquivo_ia() detecta tipo, chama Gemini para gerar estrutura
#    3. Arquivo criado/modificado, aberto automaticamente
# ============================================================

import json
import os
import re
import subprocess
import tempfile
from datetime import datetime
from pathlib import Path
from typing import Any

from rony_paths import RONY_DIR

# ── Pasta padrão para arquivos criados ───────────────────────
_PASTA_CRIADOS = Path.home() / "Documents" / "Rony"
_PASTA_CRIADOS.mkdir(parents=True, exist_ok=True)

# ── Detecção de bibliotecas ───────────────────────────────────
try:
    from docx import Document as _DocxDocument
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    DOCX_OK = True
except ImportError:
    DOCX_OK = False

try:
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter
    from openpyxl.chart import BarChart, LineChart, Reference
    XLSX_OK = True
except ImportError:
    XLSX_OK = False

try:
    from pptx import Presentation as _Pptx
    from pptx.util import Inches as PptInches, Pt as PptPt
    from pptx.dml.color import RGBColor as PptRGB
    PPTX_OK = True
except ImportError:
    PPTX_OK = False


# ═══════════════════════════════════════════════════════════════
#  UTILITÁRIOS
# ═══════════════════════════════════════════════════════════════

def _sanitizar_nome(nome: str) -> str:
    """Remove caracteres inválidos do nome de arquivo."""
    nome = re.sub(r'[<>:"/\\|?*]', "_", nome).strip(" .")
    return nome or "arquivo_rony"


def _abrir_arquivo(caminho: Path) -> None:
    """Abre o arquivo com o programa padrão do sistema."""
    try:
        os.startfile(str(caminho))
    except Exception:
        try:
            subprocess.Popen(["explorer", str(caminho)], shell=False)
        except Exception:
            pass


def _gerar_nome_arquivo(instrucao: str, extensao: str) -> Path:
    """Gera nome de arquivo baseado na instrução do usuário."""
    # Extrai palavras-chave da instrução (até 4 palavras)
    palavras = re.sub(r"[^\w\s]", "", instrucao.lower()).split()
    ignorar = {"cria", "criar", "crie", "faz", "fazer", "faça", "faca",
                "um", "uma", "uns", "umas", "de", "do", "da", "com",
                "para", "sobre", "arquivo", "planilha", "documento",
                "apresentação", "apresentacao", "make", "create", "the", "a"}
    palavras_chave = [p for p in palavras if p not in ignorar][:4]
    base = "_".join(palavras_chave) if palavras_chave else "rony_doc"
    ts = datetime.now().strftime("%Y%m%d_%H%M")
    nome = f"{_sanitizar_nome(base)}_{ts}{extensao}"
    return _PASTA_CRIADOS / nome


# ═══════════════════════════════════════════════════════════════
#  ARQUIVOS DE TEXTO (TXT / MD / CSV / JSON / Código)
# ═══════════════════════════════════════════════════════════════

def criar_txt(instrucao: str, conteudo: str, nome: str = "") -> Path:
    """Cria um arquivo de texto com o conteúdo fornecido."""
    if nome:
        if not Path(nome).suffix:
            nome += ".txt"
        caminho = _PASTA_CRIADOS / _sanitizar_nome(nome)
    else:
        caminho = _gerar_nome_arquivo(instrucao, ".txt")
    caminho.write_text(conteudo, encoding="utf-8")
    return caminho


def editar_txt(caminho: Path, novo_conteudo: str) -> Path:
    """Sobrescreve ou modifica um arquivo de texto."""
    caminho.write_text(novo_conteudo, encoding="utf-8")
    return caminho


def ler_txt(caminho: Path) -> str:
    """Lê e retorna o conteúdo de um arquivo de texto."""
    try:
        return caminho.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return caminho.read_text(encoding="latin-1")


# ═══════════════════════════════════════════════════════════════
#  WORD (.docx)
# ═══════════════════════════════════════════════════════════════

def criar_docx(instrucao: str, estrutura: dict, nome: str = "") -> Path:
    """
    Cria um arquivo Word (.docx).

    estrutura = {
        "titulo": "Relatório de Vendas",
        "subtitulo": "Janeiro 2025",
        "secoes": [
            {"titulo": "Introdução", "paragrafos": ["Texto...", "Texto..."]},
            {"titulo": "Dados", "tabela": {"cabecalho": ["A","B"], "linhas": [[1,2],[3,4]]}},
        ]
    }
    """
    if not DOCX_OK:
        raise ImportError("python-docx não instalado. Execute: pip install python-docx")

    if nome:
        if not nome.endswith(".docx"):
            nome += ".docx"
        caminho = _PASTA_CRIADOS / _sanitizar_nome(nome)
    else:
        caminho = _gerar_nome_arquivo(instrucao, ".docx")

    doc = _DocxDocument()

    # Estilo de título
    titulo = estrutura.get("titulo", "Documento")
    t = doc.add_heading(titulo, level=0)
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER

    if subtitulo := estrutura.get("subtitulo", ""):
        p = doc.add_paragraph(subtitulo)
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.runs[0].font.size = Pt(12)
        p.runs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    doc.add_paragraph()  # espaço

    # Seções
    for secao in estrutura.get("secoes", []):
        if titulo_secao := secao.get("titulo", ""):
            doc.add_heading(titulo_secao, level=1)

        for para in secao.get("paragrafos", []):
            doc.add_paragraph(str(para))

        # Tabela
        if tabela := secao.get("tabela"):
            cab = tabela.get("cabecalho", [])
            linhas = tabela.get("linhas", [])
            if cab:
                n_cols = len(cab)
                tbl = doc.add_table(rows=1 + len(linhas), cols=n_cols)
                tbl.style = "Table Grid"
                # Cabeçalho
                row_cab = tbl.rows[0]
                for i, header in enumerate(cab):
                    cell = row_cab.cells[i]
                    cell.text = str(header)
                    cell.paragraphs[0].runs[0].font.bold = True
                # Dados
                for r_idx, linha in enumerate(linhas):
                    for c_idx, valor in enumerate(linha):
                        tbl.rows[r_idx + 1].cells[c_idx].text = str(valor)

        # Lista
        for item in secao.get("lista", []):
            doc.add_paragraph(str(item), style="List Bullet")

        doc.add_paragraph()  # espaço entre seções

    # Rodapé com data
    rodape = doc.sections[0].footer
    p_rodape = rodape.paragraphs[0]
    p_rodape.text = f"Gerado pelo R.O.N.Y — {datetime.now().strftime('%d/%m/%Y %H:%M')}"
    p_rodape.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.save(str(caminho))
    return caminho


def editar_docx(caminho: Path, estrutura_adicional: dict) -> Path:
    """Adiciona conteúdo a um .docx existente."""
    if not DOCX_OK:
        raise ImportError("python-docx não instalado")
    doc = _DocxDocument(str(caminho))
    for secao in estrutura_adicional.get("secoes", []):
        if t := secao.get("titulo", ""):
            doc.add_heading(t, level=1)
        for para in secao.get("paragrafos", []):
            doc.add_paragraph(str(para))
    doc.save(str(caminho))
    return caminho


# ═══════════════════════════════════════════════════════════════
#  EXCEL (.xlsx)
# ═══════════════════════════════════════════════════════════════

def criar_xlsx(instrucao: str, estrutura: dict, nome: str = "") -> Path:
    """
    Cria planilha Excel formatada.

    estrutura = {
        "titulo": "Gastos de Janeiro",
        "abas": [
            {
                "nome": "Gastos",
                "cabecalho": ["Data", "Descrição", "Valor", "Categoria"],
                "linhas": [
                    ["01/01/2025", "Supermercado", 350.00, "Alimentação"],
                    ...
                ],
                "totais": True,       # soma automática
                "grafico": "barra"    # barra | linha | nenhum
            }
        ]
    }
    """
    if not XLSX_OK:
        raise ImportError("openpyxl não instalado")

    if nome:
        if not nome.endswith(".xlsx"):
            nome += ".xlsx"
        caminho = _PASTA_CRIADOS / _sanitizar_nome(nome)
    else:
        caminho = _gerar_nome_arquivo(instrucao, ".xlsx")

    wb = openpyxl.Workbook()
    wb.remove(wb.active)   # remove aba padrão

    # Cores
    COR_HEADER = "2196F3"    # azul
    COR_HEADER_FONT = "FFFFFF"
    COR_LINHA_PAR = "E3F2FD"
    COR_TOTAL = "1565C0"

    for aba_def in estrutura.get("abas", [{"nome": "Dados", "cabecalho": [], "linhas": []}]):
        ws = wb.create_sheet(title=aba_def.get("nome", "Dados")[:31])
        cab = aba_def.get("cabecalho", [])
        linhas = aba_def.get("linhas", [])

        linha_inicio = 1

        # Título da aba
        if titulo_aba := estrutura.get("titulo", ""):
            ws.cell(1, 1, titulo_aba).font = Font(bold=True, size=14, color="1565C0")
            ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=max(len(cab), 1))
            linha_inicio = 2
            ws.cell(2, 1, f"Gerado em {datetime.now().strftime('%d/%m/%Y %H:%M')}").font = Font(italic=True, size=9, color="888888")
            linha_inicio = 3

        # Cabeçalho
        if cab:
            row_h = linha_inicio
            for c_idx, h in enumerate(cab, 1):
                cell = ws.cell(row_h, c_idx, str(h))
                cell.font = Font(bold=True, color=COR_HEADER_FONT)
                cell.fill = PatternFill("solid", fgColor=COR_HEADER)
                cell.alignment = Alignment(horizontal="center", vertical="center")
            ws.row_dimensions[row_h].height = 22
            linha_inicio += 1

        # Dados
        col_numerica = set()
        for r_idx, linha in enumerate(linhas):
            row_r = linha_inicio + r_idx
            fill_cor = COR_LINHA_PAR if r_idx % 2 == 0 else "FFFFFF"
            for c_idx, valor in enumerate(linha, 1):
                cell = ws.cell(row_r, c_idx, valor)
                cell.fill = PatternFill("solid", fgColor=fill_cor)
                cell.alignment = Alignment(vertical="center")
                if isinstance(valor, (int, float)):
                    cell.number_format = '#,##0.00' if isinstance(valor, float) else '#,##0'
                    col_numerica.add(c_idx)

        # Linha de totais
        if aba_def.get("totais") and linhas and cab:
            row_tot = linha_inicio + len(linhas)
            ws.cell(row_tot, 1, "TOTAL").font = Font(bold=True, color="FFFFFF")
            ws.cell(row_tot, 1).fill = PatternFill("solid", fgColor=COR_TOTAL)
            for c_idx in col_numerica:
                inicio_formula = linha_inicio
                fim_formula = linha_inicio + len(linhas) - 1
                col_letra = get_column_letter(c_idx)
                formula = f"=SUM({col_letra}{inicio_formula}:{col_letra}{fim_formula})"
                cell = ws.cell(row_tot, c_idx, formula)
                cell.font = Font(bold=True, color="FFFFFF")
                cell.fill = PatternFill("solid", fgColor=COR_TOTAL)
                cell.number_format = '#,##0.00'

        # Auto largura das colunas
        for col in ws.columns:
            max_w = 0
            for cell in col:
                try:
                    if cell.value:
                        max_w = max(max_w, len(str(cell.value)))
                except Exception:
                    pass
            ws.column_dimensions[get_column_letter(col[0].column)].width = min(max_w + 4, 40)

        # Gráfico
        tipo_grafico = aba_def.get("grafico", "nenhum")
        if tipo_grafico in ("barra", "linha") and linhas and len(cab) >= 2:
            try:
                row_data_start = linha_inicio if not aba_def.get("cabecalho") else linha_inicio
                row_data_end = row_data_start + len(linhas) - 1
                # Usa primeira coluna como categorias, segunda como dados
                cats = Reference(ws, min_col=1, min_row=row_data_start, max_row=row_data_end)
                data = Reference(ws, min_col=2, min_row=row_data_start - 1, max_row=row_data_end)

                chart = BarChart() if tipo_grafico == "barra" else LineChart()
                chart.add_data(data, titles_from_data=True)
                chart.set_categories(cats)
                chart.title = estrutura.get("titulo", "Gráfico")
                chart.style = 10
                chart.width = 20
                chart.height = 12

                # Posiciona após os dados
                col_grafico = get_column_letter(len(cab) + 2)
                ws.add_chart(chart, f"{col_grafico}{linha_inicio}")
            except Exception as _e:
                print(f"[EXCEL] Grafico nao criado: {_e}")

    wb.save(str(caminho))
    return caminho


# ═══════════════════════════════════════════════════════════════
#  POWERPOINT (.pptx)
# ═══════════════════════════════════════════════════════════════

def criar_pptx(instrucao: str, estrutura: dict, nome: str = "") -> Path:
    """
    Cria apresentação PowerPoint.

    estrutura = {
        "titulo": "Apresentação de Vendas",
        "autor": "Robert",
        "slides": [
            {"tipo": "titulo", "titulo": "Vendas 2025", "subtitulo": "Resultados Q1"},
            {"tipo": "conteudo", "titulo": "Pontos Principais", "bullets": ["Ponto 1", "Ponto 2"]},
            {"tipo": "tabela", "titulo": "Dados", "cabecalho": ["A","B"], "linhas": [[1,2]]},
            {"tipo": "encerramento", "mensagem": "Obrigado!"},
        ]
    }
    """
    if not PPTX_OK:
        raise ImportError("python-pptx não instalado")

    if nome:
        if not nome.endswith(".pptx"):
            nome += ".pptx"
        caminho = _PASTA_CRIADOS / _sanitizar_nome(nome)
    else:
        caminho = _gerar_nome_arquivo(instrucao, ".pptx")

    prs = _Pptx()
    prs.slide_width  = PptInches(13.33)
    prs.slide_height = PptInches(7.5)

    layouts = prs.slide_layouts

    for slide_def in estrutura.get("slides", []):
        tipo = slide_def.get("tipo", "conteudo")

        if tipo == "titulo":
            slide = prs.slides.add_slide(layouts[0])  # título
            if slide.shapes.title:
                slide.shapes.title.text = slide_def.get("titulo", "")
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = slide_def.get("subtitulo", "")

        elif tipo in ("conteudo", "bullets"):
            slide = prs.slides.add_slide(layouts[1])  # título + conteúdo
            if slide.shapes.title:
                slide.shapes.title.text = slide_def.get("titulo", "")
            bullets = slide_def.get("bullets", slide_def.get("paragrafos", []))
            if bullets and len(slide.placeholders) > 1:
                tf = slide.placeholders[1].text_frame
                tf.clear()
                for i, bullet in enumerate(bullets):
                    if i == 0:
                        tf.paragraphs[0].text = str(bullet)
                    else:
                        p = tf.add_paragraph()
                        p.text = str(bullet)

        elif tipo == "tabela":
            slide = prs.slides.add_slide(layouts[5])  # em branco
            # Título manual
            from pptx.util import Emu
            tx_box = slide.shapes.add_textbox(PptInches(0.5), PptInches(0.3),
                                              PptInches(12), PptInches(0.8))
            tx_box.text_frame.text = slide_def.get("titulo", "")
            tx_box.text_frame.paragraphs[0].font.bold = True
            tx_box.text_frame.paragraphs[0].font.size = PptPt(28)

            cab  = slide_def.get("cabecalho", [])
            rows = slide_def.get("linhas", [])
            if cab and rows:
                tbl = slide.shapes.add_table(
                    1 + len(rows), len(cab),
                    PptInches(0.5), PptInches(1.5),
                    PptInches(12), PptInches(min(4.5, 0.4 * (len(rows) + 1)))
                ).table
                for c_idx, h in enumerate(cab):
                    tbl.cell(0, c_idx).text = str(h)
                for r_idx, linha in enumerate(rows):
                    for c_idx, v in enumerate(linha):
                        tbl.cell(r_idx + 1, c_idx).text = str(v)

        elif tipo == "encerramento":
            slide = prs.slides.add_slide(layouts[6])  # em branco centrado
            tx = slide.shapes.add_textbox(PptInches(1), PptInches(2.5),
                                           PptInches(11), PptInches(2))
            tf = tx.text_frame
            tf.text = slide_def.get("mensagem", "Obrigado!")
            tf.paragraphs[0].font.size = PptPt(44)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].alignment = 2  # center

        else:
            slide = prs.slides.add_slide(layouts[1])
            if slide.shapes.title:
                slide.shapes.title.text = slide_def.get("titulo", tipo)

    prs.save(str(caminho))
    return caminho


# ═══════════════════════════════════════════════════════════════
#  HTML DASHBOARD
# ═══════════════════════════════════════════════════════════════

def criar_html_dashboard(instrucao: str, estrutura: dict, nome: str = "") -> Path:
    """
    Gera dashboard HTML completo com Chart.js.

    estrutura = {
        "titulo": "Dashboard de Vendas",
        "cards": [
            {"label": "Total", "valor": "R$ 45.200", "icone": "💰", "cor": "#2196F3"},
        ],
        "graficos": [
            {
                "tipo": "bar",      # bar | line | pie | doughnut
                "titulo": "Vendas por Mês",
                "labels": ["Jan", "Fev", "Mar"],
                "datasets": [{"label": "2025", "data": [10, 20, 15]}]
            }
        ],
        "tabelas": [
            {
                "titulo": "Top Produtos",
                "cabecalho": ["Produto", "Qtd", "Valor"],
                "linhas": [["Produto A", 100, "R$ 5.000"]]
            }
        ]
    }
    """
    if nome:
        if not nome.endswith(".html"):
            nome += ".html"
        caminho = _PASTA_CRIADOS / _sanitizar_nome(nome)
    else:
        caminho = _gerar_nome_arquivo(instrucao, ".html")

    titulo = estrutura.get("titulo", "Dashboard R.O.N.Y")
    cards  = estrutura.get("cards", [])
    graficos = estrutura.get("graficos", [])
    tabelas  = estrutura.get("tabelas", [])

    # ── HTML dos cards ─────────────────────────────────────────
    html_cards = ""
    for card in cards:
        cor = card.get("cor", "#2196F3")
        html_cards += f"""
        <div class="card" style="border-left:4px solid {cor}">
            <div class="card-icon">{card.get("icone","📊")}</div>
            <div class="card-body">
                <div class="card-value">{card.get("valor","—")}</div>
                <div class="card-label">{card.get("label","")}</div>
            </div>
        </div>"""

    # ── HTML dos gráficos ──────────────────────────────────────
    html_graficos = ""
    js_graficos   = ""
    CORES = ["#2196F3","#4CAF50","#FF9800","#9C27B0","#F44336","#00BCD4","#FF5722","#607D8B"]

    for i, g in enumerate(graficos):
        gid = f"chart_{i}"
        tipo = g.get("tipo", "bar")
        labels = json.dumps(g.get("labels", []))
        datasets_def = g.get("datasets", [])
        datasets_js = []
        for j, ds in enumerate(datasets_def):
            cor = ds.get("cor", CORES[j % len(CORES)])
            datasets_js.append(json.dumps({
                "label": ds.get("label", f"Série {j+1}"),
                "data": ds.get("data", []),
                "backgroundColor": cor + "CC",
                "borderColor": cor,
                "borderWidth": 2,
                "fill": tipo == "line",
                "tension": 0.4,
            }))
        datasets_str = "[" + ",".join(datasets_js) + "]"

        html_graficos += f"""
        <div class="chart-box">
            <h3>{g.get("titulo","")}</h3>
            <canvas id="{gid}"></canvas>
        </div>"""

        js_graficos += f"""
        new Chart(document.getElementById('{gid}'), {{
            type: '{tipo}',
            data: {{ labels: {labels}, datasets: {datasets_str} }},
            options: {{
                responsive: true, maintainAspectRatio: true,
                plugins: {{ legend: {{ position: 'bottom' }} }},
                scales: {{ y: {{ beginAtZero: true }} }}
            }}
        }});"""

    # ── HTML das tabelas ───────────────────────────────────────
    html_tabelas = ""
    for tab in tabelas:
        cab = tab.get("cabecalho", [])
        html_th = "".join(f"<th>{h}</th>" for h in cab)
        html_linhas = ""
        for linha in tab.get("linhas", []):
            html_linhas += "<tr>" + "".join(f"<td>{v}</td>" for v in linha) + "</tr>"
        html_tabelas += f"""
        <div class="table-box">
            <h3>{tab.get("titulo","Tabela")}</h3>
            <table><thead><tr>{html_th}</tr></thead><tbody>{html_linhas}</tbody></table>
        </div>"""

    # ── HTML completo ──────────────────────────────────────────
    agora = datetime.now().strftime("%d/%m/%Y %H:%M")
    html = f"""<!DOCTYPE html>
<html lang="pt-BR">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>{titulo}</title>
<script src="https://cdn.jsdelivr.net/npm/chart.js"></script>
<style>
  * {{ box-sizing: border-box; margin: 0; padding: 0; }}
  body {{ font-family: 'Segoe UI', Arial, sans-serif; background: #f0f2f5; color: #333; }}
  header {{ background: linear-gradient(135deg, #1a237e, #1565c0); color: white;
            padding: 24px 32px; display: flex; align-items: center; justify-content: space-between; }}
  header h1 {{ font-size: 1.8rem; font-weight: 700; }}
  header .ts {{ font-size: 0.85rem; opacity: 0.75; }}
  .container {{ max-width: 1400px; margin: 0 auto; padding: 24px; }}
  .cards {{ display: grid; grid-template-columns: repeat(auto-fit, minmax(200px, 1fr));
            gap: 16px; margin-bottom: 24px; }}
  .card {{ background: white; border-radius: 12px; padding: 20px 24px;
           display: flex; align-items: center; gap: 16px;
           box-shadow: 0 2px 8px rgba(0,0,0,0.08); }}
  .card-icon {{ font-size: 2rem; }}
  .card-value {{ font-size: 1.6rem; font-weight: 700; color: #1a237e; }}
  .card-label {{ font-size: 0.82rem; color: #888; margin-top: 4px; }}
  .charts {{ display: grid; grid-template-columns: repeat(auto-fit, minmax(400px, 1fr));
             gap: 20px; margin-bottom: 24px; }}
  .chart-box {{ background: white; border-radius: 12px; padding: 20px;
                box-shadow: 0 2px 8px rgba(0,0,0,0.08); }}
  .chart-box h3 {{ margin-bottom: 16px; color: #1565c0; font-size: 1rem; }}
  .tables {{ display: grid; grid-template-columns: repeat(auto-fit, minmax(400px, 1fr)); gap: 20px; }}
  .table-box {{ background: white; border-radius: 12px; padding: 20px;
                box-shadow: 0 2px 8px rgba(0,0,0,0.08); overflow-x: auto; }}
  .table-box h3 {{ margin-bottom: 16px; color: #1565c0; font-size: 1rem; }}
  table {{ width: 100%; border-collapse: collapse; font-size: 0.9rem; }}
  th {{ background: #1565c0; color: white; padding: 10px 12px; text-align: left; }}
  td {{ padding: 9px 12px; border-bottom: 1px solid #eee; }}
  tr:nth-child(even) td {{ background: #f5f8ff; }}
  tr:hover td {{ background: #e3f2fd; }}
  footer {{ text-align: center; padding: 20px; color: #aaa; font-size: 0.8rem; }}
</style>
</head>
<body>
<header>
  <h1>📊 {titulo}</h1>
  <span class="ts">Gerado pelo R.O.N.Y · {agora}</span>
</header>
<div class="container">
  <div class="cards">{html_cards}</div>
  <div class="charts">{html_graficos}</div>
  <div class="tables">{html_tabelas}</div>
</div>
<footer>R.O.N.Y — Assistente Pessoal · {agora}</footer>
<script>
{js_graficos}
</script>
</body>
</html>"""

    caminho.write_text(html, encoding="utf-8")
    return caminho


# ═══════════════════════════════════════════════════════════════
#  PONTO DE ENTRADA PRINCIPAL — IA-powered
# ═══════════════════════════════════════════════════════════════

async def editar_arquivo_ia(
    instrucao: str,
    client_gemini=None,
    lingua: str = "pt",
    abrir: bool = True,
) -> str:
    """
    Interpreta a instrução do usuário, usa Gemini para gerar a estrutura
    do arquivo e cria o arquivo correspondente.

    Retorna uma mensagem de confirmação para ser falada pelo Rony.
    """
    import asyncio

    instrucao_lower = instrucao.lower()

    # ── Detecta tipo de arquivo ───────────────────────────────
    if any(k in instrucao_lower for k in ["excel", "planilha", "xlsx", "spreadsheet", "tabela", "dashboard financeiro"]):
        tipo = "xlsx"
    elif any(k in instrucao_lower for k in ["dashboard", "html", "página", "pagina", "site", "relatorio html"]):
        tipo = "html"
    elif any(k in instrucao_lower for k in ["word", "docx", "documento", "relatorio", "relatório", "carta", "contrato", "redacao", "redação"]):
        tipo = "docx"
    elif any(k in instrucao_lower for k in ["apresentação", "apresentacao", "slide", "pptx", "powerpoint", "pitch"]):
        tipo = "pptx"
    else:
        tipo = "txt"

    # ── Prompt para Gemini gerar a estrutura ──────────────────
    prompts = {
        "xlsx": f"""Você é um assistente que cria planilhas Excel.
O usuário pediu: "{instrucao}"

Responda APENAS com um JSON válido (sem markdown, sem explicações) com esta estrutura:
{{
  "titulo": "Título da planilha",
  "abas": [
    {{
      "nome": "Nome da aba",
      "cabecalho": ["Coluna1", "Coluna2", "Coluna3"],
      "linhas": [
        ["valor1", "valor2", 100.00],
        ["valor3", "valor4", 200.00]
      ],
      "totais": true,
      "grafico": "barra"
    }}
  ]
}}
Gere dados realistas e relevantes para o pedido. Inclua pelo menos 6 linhas de dados.""",

        "html": f"""Você é um assistente que cria dashboards HTML.
O usuário pediu: "{instrucao}"

Responda APENAS com um JSON válido (sem markdown) com esta estrutura:
{{
  "titulo": "Título do Dashboard",
  "cards": [
    {{"label": "Label", "valor": "Valor", "icone": "emoji", "cor": "#hexcor"}}
  ],
  "graficos": [
    {{
      "tipo": "bar",
      "titulo": "Título",
      "labels": ["Jan","Fev","Mar"],
      "datasets": [{{"label": "Série", "data": [10,20,15]}}]
    }}
  ],
  "tabelas": [
    {{
      "titulo": "Título",
      "cabecalho": ["Col1","Col2"],
      "linhas": [["v1","v2"]]
    }}
  ]
}}
Gere dados realistas e relevantes. Inclua pelo menos 2 gráficos e 4 cards.""",

        "docx": f"""Você é um assistente que cria documentos Word.
O usuário pediu: "{instrucao}"

Responda APENAS com um JSON válido (sem markdown) com esta estrutura:
{{
  "titulo": "Título do Documento",
  "subtitulo": "Subtítulo opcional",
  "secoes": [
    {{
      "titulo": "Nome da Seção",
      "paragrafos": ["Parágrafo 1...", "Parágrafo 2..."],
      "lista": ["Item 1", "Item 2"],
      "tabela": {{
        "cabecalho": ["Col1","Col2"],
        "linhas": [["v1","v2"]]
      }}
    }}
  ]
}}
Gere conteúdo completo, profissional e relevante para o pedido.""",

        "pptx": f"""Você é um assistente que cria apresentações PowerPoint.
O usuário pediu: "{instrucao}"

Responda APENAS com um JSON válido (sem markdown) com esta estrutura:
{{
  "titulo": "Título da Apresentação",
  "autor": "Rony",
  "slides": [
    {{"tipo": "titulo", "titulo": "Título Principal", "subtitulo": "Subtítulo"}},
    {{"tipo": "conteudo", "titulo": "Ponto 1", "bullets": ["Item A","Item B","Item C"]}},
    {{"tipo": "tabela", "titulo": "Dados", "cabecalho": ["A","B"], "linhas": [["v1","v2"]]}},
    {{"tipo": "encerramento", "mensagem": "Obrigado!"}}
  ]
}}
Crie uma apresentação completa com pelo menos 5 slides relevantes para o pedido.""",

        "txt": f"""O usuário pediu: "{instrucao}"
Gere o conteúdo de texto completo para atender ao pedido.
Seja direto, organizado e use formatação com markdown se fizer sentido.
Responda APENAS com o texto do arquivo, sem explicações adicionais.""",
    }

    estrutura = None
    conteudo_txt = ""

    if client_gemini:
        try:
            import google.genai as genai

            resp = await asyncio.get_event_loop().run_in_executor(
                None,
                lambda: client_gemini.models.generate_content(
                    model="gemini-2.0-flash",
                    contents=prompts[tipo],
                ),
            )
            resposta_raw = resp.text.strip()

            if tipo != "txt":
                # Remove blocos markdown se presentes
                resposta_raw = re.sub(r"```json\s*", "", resposta_raw)
                resposta_raw = re.sub(r"```\s*$", "", resposta_raw, flags=re.MULTILINE)
                try:
                    estrutura = json.loads(resposta_raw)
                except json.JSONDecodeError:
                    # Tenta encontrar JSON na resposta
                    match = re.search(r"\{[\s\S]+\}", resposta_raw)
                    if match:
                        estrutura = json.loads(match.group())
            else:
                conteudo_txt = resposta_raw

        except Exception as e:
            print(f"[FILE_EDITOR] Gemini falhou: {e} — usando estrutura padrão")

    # ── Fallback sem IA ───────────────────────────────────────
    if estrutura is None and tipo != "txt":
        estrutura = _estrutura_padrao(tipo, instrucao)
    if not conteudo_txt and tipo == "txt":
        conteudo_txt = f"# {instrucao}\n\nConteúdo gerado pelo R.O.N.Y em {datetime.now().strftime('%d/%m/%Y %H:%M')}.\n"

    # ── Cria o arquivo ────────────────────────────────────────
    try:
        if tipo == "xlsx":
            caminho = criar_xlsx(instrucao, estrutura)
        elif tipo == "html":
            caminho = criar_html_dashboard(instrucao, estrutura)
        elif tipo == "docx":
            caminho = criar_docx(instrucao, estrutura)
        elif tipo == "pptx":
            caminho = criar_pptx(instrucao, estrutura)
        else:
            caminho = criar_txt(instrucao, conteudo_txt)

        if abrir:
            _abrir_arquivo(caminho)

        msgs = {
            "pt": f"Arquivo criado em {caminho.name}. Abrindo agora.",
            "en": f"File {caminho.name} created and opened.",
            "fr": f"Fichier {caminho.name} créé et ouvert.",
            "es": f"Archivo {caminho.name} creado y abierto.",
        }
        print(f"[FILE_EDITOR] Criado: {caminho}")
        return msgs.get(lingua, msgs["pt"])

    except Exception as e:
        print(f"[FILE_EDITOR] Erro ao criar arquivo: {e}")
        msgs = {
            "pt": f"Não consegui criar o arquivo. Erro: {e}",
            "en": f"Could not create the file. Error: {e}",
        }
        return msgs.get(lingua, msgs["pt"])


def _estrutura_padrao(tipo: str, instrucao: str) -> dict:
    """Estrutura mínima quando Gemini não está disponível."""
    nome = instrucao[:40]
    if tipo == "xlsx":
        return {
            "titulo": nome,
            "abas": [{"nome": "Dados", "cabecalho": ["Item", "Valor", "Data"],
                       "linhas": [["Exemplo", 100.0, datetime.now().strftime("%d/%m/%Y")]],
                       "totais": True, "grafico": "barra"}]
        }
    if tipo == "html":
        return {
            "titulo": nome,
            "cards": [{"label": "Total", "valor": "—", "icone": "📊", "cor": "#2196F3"}],
            "graficos": [{"tipo": "bar", "titulo": "Dados", "labels": ["A","B","C"],
                          "datasets": [{"label": "Série 1", "data": [10,20,15]}]}],
            "tabelas": []
        }
    if tipo == "docx":
        return {
            "titulo": nome,
            "secoes": [{"titulo": "Conteúdo", "paragrafos": ["Documento gerado pelo R.O.N.Y."]}]
        }
    if tipo == "pptx":
        return {
            "titulo": nome,
            "slides": [
                {"tipo": "titulo", "titulo": nome, "subtitulo": "R.O.N.Y"},
                {"tipo": "conteudo", "titulo": "Conteúdo", "bullets": ["Slide gerado automaticamente."]},
                {"tipo": "encerramento", "mensagem": "Obrigado!"},
            ]
        }
    return {}
